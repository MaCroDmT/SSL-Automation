from pptx import Presentation
from pptx.util import Inches, Pt
import os
from pathlib import Path

# --- HELPER: Get Downloads Folder Path ---
def get_downloads_path():
    """Returns the absolute path to the user's Downloads folder."""
    return str(Path.home() / "Downloads")

class PresentationGenerator:
    def __init__(self, template_path=None):
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()

    def create_title_slide(self, title_text, subtitle_text):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        try:
            slide.placeholders[0].text = title_text
            slide.placeholders[1].text = subtitle_text
        except IndexError:
            pass

    def create_content_slide(self, title_text, content_text):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        try:
            tf = slide.placeholders[1].text_frame
            tf.text = content_text 
        except IndexError:
            pass

    def create_image_slide(self, title_text, image_path, description):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        
        # --- IMAGE DEBUGGING & INSERTION ---
        print(f"Looking for image at: {image_path}")
        
        if image_path and os.path.exists(image_path):
            left = Inches(2)
            top = Inches(2)
            height = Inches(4.5) 
            slide.shapes.add_picture(image_path, left, top, height=height)
            print(" -> Image found and inserted.")
        else:
            print(f" -> ERROR: Image NOT found. Check filename or path.")
        
        txBox = slide.shapes.add_textbox(Inches(3), Inches(6.6), Inches(8), Inches(1))
        txBox.text_frame.text = description

    def save(self, filename):
        full_path = os.path.join(get_downloads_path(), filename)
        try:
            self.prs.save(full_path)
            print("-" * 40)
            print(f"SUCCESS: Presentation saved to: {full_path}")
            print("-" * 40)
        except PermissionError:
            print(f"ERROR: Could not write to {full_path}. Close the file if open.")

# --- USER INPUT SECTION (UPDATED) ---

# 1. Define the correct path to your image
# This combines "C:\Users\User\Downloads" with "F16.jpg"
image_location = os.path.join(get_downloads_path(), "F16.jpg")

user_request = [
    {
        "type": "title_slide",
        "title": "F-16 Fighting Falcon",
        "subtitle": "Overview of the Multirole Fighter"
    },
    {
        "type": "bullet_slide",
        "title": "Key Features",
        "content": "• High maneuverability\n• Mach 2+ speeds\n• Advanced avionics systems"
    },
    {
        "type": "image_slide",
        "title": "Visual Profile",
        "image": image_location,  # <--- USING THE CORRECT PATH HERE
        "caption": "The F-16 in a high-G maneuver."
    }
]

# --- MAIN EXECUTION ---
if __name__ == "__main__":
    ppt_bot = PresentationGenerator() 

    for slide_data in user_request:
        if slide_data['type'] == 'title_slide':
            ppt_bot.create_title_slide(slide_data['title'], slide_data['subtitle'])
        elif slide_data['type'] == 'bullet_slide':
            ppt_bot.create_content_slide(slide_data['title'], slide_data['content'])
        elif slide_data['type'] == 'image_slide':
            ppt_bot.create_image_slide(slide_data['title'], slide_data.get('image'), slide_data.get('caption'))

    ppt_bot.save("Generated_Presentation.pptx")